import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_DARK = RGBColor(10, 22, 40)       # #0A1628
    CARD_BG = RGBColor(17, 34, 64)       # #112240
    ACCENT_BLUE = RGBColor(26, 115, 232)  # #1A73E8
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(136, 146, 176)
    ACCENT_TEAL = RGBColor(0, 201, 167)
    
    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text=""):
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = "BLUESTOCK ANALYTICS  |  " + title_text.upper()
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_BLUE
        
        if subtitle_text:
            p1 = tf.add_paragraph()
            p1.text = subtitle_text
            p1.font.size = Pt(12)
            p1.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, title, content_list, accent_color=ACCENT_BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = accent_color
        card.line.width = Pt(1.5)
        
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = "• " + item
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = TEXT_WHITE
            p_item.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)
    
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "BLUESTOCK MUTUAL FUND ANALYTICS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    p2 = tf1.add_paragraph()
    p2.text = "Comprehensive Industry Capstone Report & Performance Dashboard"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(12)
    
    p3 = tf1.add_paragraph()
    p3.text = "ETL Pipeline • Quantitative Risk Modeling • Power BI Visual Analytics • Investor Behavior"
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT_TEAL
    p3.space_before = Pt(24)

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Strategic Objectives
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Problem Statement & Objectives", "Aligning Quantitative Analytics with Retail Investor Growth")
    
    add_card(s2, 0.8, 1.6, 5.6, 5.2, "Business Problem", [
        "Rapid retail mutual fund adoption requires automated data pipelines.",
        "Lack of unified performance & risk scorecards across AMCs.",
        "High investor churn and gap between SIP registration and fulfillment.",
        "Need for data-driven risk management (Sharpe, Sortino, VaR/CVaR)."
    ], ACCENT_BLUE)
    
    add_card(s2, 6.8, 1.6, 5.7, 5.2, "Strategic Objectives", [
        "Build automated SQLite ETL pipeline importing 10 raw datasets.",
        "Compute 1yr/3yr/5yr CAGR, Sharpe, Sortino, Alpha, Beta, & Max Drawdown.",
        "Construct 0-100 Composite Fund Scorecard across 40 schemes.",
        "Design 4-page interactive Power BI dashboard & risk recommender."
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 3: Data Sources & Dataset Overview
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Data Sources & AMFI Dataset Overview", "Comprehensive Financial & Demographic Data Integration")
    
    add_card(s3, 0.8, 1.6, 3.6, 5.2, "Fund Metadata & NAV", [
        "40 AMFI Scheme Master (`dim_fund`).",
        "46,000 daily NAV observations (2022–2026).",
        "Expense ratios & benchmark mappings."
    ])
    
    add_card(s3, 4.8, 1.6, 3.6, 5.2, "Industry Metrics", [
        "Industry AUM history by AMC (`fact_aum`).",
        "Monthly SIP inflows (`fact_monthly_sip`).",
        "Category net inflows & Folio counts."
    ], ACCENT_TEAL)
    
    add_card(s3, 8.8, 1.6, 3.6, 5.2, "Investor & Holdings", [
        "32,778 anonymized investor transactions.",
        "Demographic state & city tier splits.",
        "322 stock-level portfolio holdings."
    ], ACCENT_BLUE)

    # -------------------------------------------------------------
    # SLIDE 4: Architecture & ETL Design
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "ETL Architecture & Data Pipeline Design", "Star-Schema Warehouse & Automated Data Validation")
    
    add_card(s4, 0.8, 1.6, 3.6, 5.2, "1. Ingestion Layer", [
        "Automated extraction from CSV sources.",
        "Live NAV API fetch integration.",
        "Data type standardization & parsing."
    ])
    
    add_card(s4, 4.8, 1.6, 3.6, 5.2, "2. Quality & Clean", [
        "Forward-filling non-trading NAV gaps.",
        "Expense ratio validation (0.1%-2.5%).",
        "Positive transaction value validation."
    ], ACCENT_TEAL)
    
    add_card(s4, 8.8, 1.6, 3.6, 5.2, "3. Star-Schema Storage", [
        "`bluestock_mf.db` SQLite database.",
        "Foreign key constraints on amfi_code & date.",
        "Optimized indexes for fast analytics."
    ], ACCENT_BLUE)

    # -------------------------------------------------------------
    # SLIDE 5: EDA Highlights 1 — AUM Growth & SIP Trajectory
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "EDA Highlights: AUM & SIP Growth Trajectory", "Surging Retail Inflows & Industry Expansion")
    
    add_card(s5, 0.8, 1.6, 5.6, 5.2, "Industry AUM Expansion", [
        "Total Industry AUM reached ₹81 Lakh Crore by 2026.",
        "SBI MF, ICICI Prudential, and HDFC MF lead AUM share.",
        "Consistent positive quarterly compounding despite market swings."
    ], ACCENT_BLUE)
    
    add_card(s5, 6.8, 1.6, 5.7, 5.2, "Monthly SIP Record Inflows", [
        "Monthly SIP inflows scaled from ₹11,500 Cr to >₹31,000 Cr.",
        "YoY SIP growth exceeds 35% annually.",
        "Active SIP accounts expanded to over 7.2 Crore."
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 6: EDA Highlights 2 — Folio Trends & Investor Demographics
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "EDA Highlights: Investor Demographics & Folios", "Geographic & Generational Asset Allocations")
    
    add_card(s6, 0.8, 1.6, 5.6, 5.2, "Folio Expansion Trends", [
        "Total industry folios surpassed 26.12 Crore.",
        "Equity category accounts for >70% of total folios.",
        "Passive ETF and Index folios expanded at 45% CAGR."
    ], ACCENT_BLUE)
    
    add_card(s6, 6.8, 1.6, 5.7, 5.2, "Investor Demographic Insights", [
        "Top investment states: Maharashtra, Karnataka, Gujarat.",
        "Age group 25–44 accounts for 62% of all SIP transactions.",
        "Tier-2 & Tier-3 cities contribute 41% of new SIP registrations."
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 7: Performance Metrics 1 — Risk-Adjusted Returns & Scorecard
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Performance & Fund Scorecard Model", "Multi-Factor Composite Ranking Engine (0–100)")
    
    add_card(s7, 0.8, 1.6, 5.6, 5.2, "Scorecard Weighting Formula", [
        "30% × 3-Year CAGR Return Rank",
        "25% × Sharpe Ratio Rank (Rf = 6.5%)",
        "20% × Alpha Rank (OLS vs Nifty 100)",
        "15% × Expense Ratio Rank (Inverse / Lower is better)",
        "10% × Max Drawdown Rank (Inverse / Smaller loss is better)"
    ], ACCENT_BLUE)
    
    add_card(s7, 6.8, 1.6, 5.7, 5.2, "Top 5 Scorecard Leaders", [
        "#1 Mirae Asset Large Cap Fund — Score: 85.1 (3yr CAGR 18.1%, Sharpe 0.94)",
        "#2 ICICI Pru Midcap Fund — Score: 83.1 (3yr CAGR 22.9%, Sharpe 1.09)",
        "#3 Kotak Flexicap Fund — Score: 80.7 (3yr CAGR 19.4%, Sharpe 0.97)",
        "#4 HDFC Mid-Cap Opportunities — Score: 80.1 (3yr CAGR 21.8%, Sharpe 1.04)",
        "#5 ICICI Pru Bluechip Fund — Score: 77.7 (3yr CAGR 17.5%, Sharpe 0.89)"
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 8: Performance Metrics 2 — OLS Regression & Tail Risk
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "OLS Regression & Tail Risk (VaR / CVaR)", "Systematic Risk (Beta), Excess Return (Alpha), & Tail Downside")
    
    add_card(s8, 0.8, 1.6, 5.6, 5.2, "OLS Regression vs NIFTY 100", [
        "Single-index market model using scipy.stats.linregress.",
        "Beta values range from 0.22 (Gilt funds) to 1.04 (Small Cap).",
        "Top Alpha generators: ICICI Midcap (2.14%) & Kotak Flexicap (2.15%).",
        "Results exported to `alpha_beta.csv`."
    ], ACCENT_BLUE)
    
    add_card(s8, 6.8, 1.6, 5.7, 5.2, "Historical VaR (95%) & CVaR (95%)", [
        "VaR 95% measures 5th percentile worst daily return.",
        "CVaR (Expected Shortfall) measures average loss below VaR.",
        "Small Cap schemes display highest VaR (-2.61%) and CVaR (-3.25%).",
        "Results exported to `var_cvar_report.csv`."
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 9: Dashboard Screenshots 1 — Industry Overview & Performance
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "Power BI Dashboard: Industry & Performance", "Interactive Executive Visual Reports (Pages 1 & 2)")
    
    img_p1 = 'Dashboard/Page1_Industry_Overview.png'
    if os.path.exists(img_p1):
        s9.shapes.add_picture(img_p1, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
        
    img_p2 = 'Dashboard/Page2_Fund_Performance.png'
    if os.path.exists(img_p2):
        s9.shapes.add_picture(img_p2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))

    # -------------------------------------------------------------
    # SLIDE 10: Dashboard Screenshots 2 — Investor Analytics & Trends
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10)
    add_header(s10, "Power BI Dashboard: Investor & Market Trends", "Demographic Insights & Dual-Axis SIP Correlation (Pages 3 & 4)")
    
    img_p3 = 'Dashboard/Page3_Investor_Analytics.png'
    if os.path.exists(img_p3):
        s10.shapes.add_picture(img_p3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
        
    img_p4 = 'Dashboard/Page4_SIP_Market_Trends.png'
    if os.path.exists(img_p4):
        s10.shapes.add_picture(img_p4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))

    # -------------------------------------------------------------
    # SLIDE 11: Strategic Recommendations
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11)
    add_header(s11, "Strategic Recommendations & Insights", "Actionable Directives for AMCs & Wealth Managers")
    
    add_card(s11, 0.8, 1.6, 5.6, 5.2, "Product & Retention Directives", [
        "Address SIP Gap Latency: 97.8% of 6+ SIP investors show >35 day gaps.",
        "Promote Flexi-Cap & Mid-Cap Funds: Highest Sharpe ratio risk efficiency.",
        "Lower Expense Ratios on Direct Plans to boost scorecard competitiveness."
    ], ACCENT_BLUE)
    
    add_card(s11, 6.8, 1.6, 5.7, 5.2, "Marketing & Risk Management", [
        "Expand Tier-2 & Tier-3 SIP Outreach (41% growth driver).",
        "Incorporate VaR/CVaR Tail Risk Warnings on Small Cap offerings.",
        "Deploy Automated Risk Recommender CLI (`recommender.py`)."
    ], ACCENT_TEAL)

    # -------------------------------------------------------------
    # SLIDE 12: Thank You / Q&A
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12)
    
    tb12 = s12.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    
    p = tf12.paragraphs[0]
    p.text = "THANK YOU"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    p2 = tf12.add_paragraph()
    p2.text = "Bluestock Mutual Fund Analytics Capstone Project"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(16)
    
    p3 = tf12.add_paragraph()
    p3.text = "Deliverables: Final_Report.pdf | Bluestock_MF_Presentation.pptx | Performance_Analytics.ipynb | Advanced_Analytics.ipynb"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(12)
    p3.font.color.rgb = ACCENT_TEAL
    p3.space_before = Pt(24)

    # Save presentation
    os.makedirs('Reports', exist_ok=True)
    prs.save('Bluestock_MF_Presentation.pptx')
    prs.save('Reports/Bluestock_MF_Presentation.pptx')
    print("Saved Bluestock_MF_Presentation.pptx (12 slides) successfully.")

if __name__ == '__main__':
    create_presentation()
